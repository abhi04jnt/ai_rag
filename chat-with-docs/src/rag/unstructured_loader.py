"""
Unified loader for all unstructured data formats.

Handles text extraction and image extraction from:
- Documents: PDF (with images), DOCX, DOC, ODT, RTF
- Presentations: PPTX, PPT
- Spreadsheets: XLSX, XLS, CSV, TSV
- Web & Markup: HTML, XML
- Email: EML, MSG
- Books: EPUB
- Text & Data: TXT, MD, JSON, YAML

All format handling is consolidated in this single module.
"""
from pathlib import Path
from typing import Optional
import logging

logger = logging.getLogger(__name__)

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    # Documents
    '.pdf', '.doc', '.docx', '.odt', '.rtf',
    # Presentations
    '.ppt', '.pptx',
    # Spreadsheets
    '.xls', '.xlsx', '.csv', '.tsv',
    # Web & Markup
    '.html', '.htm', '.xml',
    # Email
    '.eml', '.msg',
    # Text & Markdown
    '.txt', '.md', '.markdown', '.rst', '.org',
    # Data formats
    '.json', '.ndjson', '.yaml', '.yml',
    # Books
    '.epub',
}

class UnstructuredDocumentLoader:
    """Load and extract text from various unstructured file formats"""
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check if required libraries are available"""
        try:
            import unstructured
            self.has_unstructured = True
        except ImportError:
            self.has_unstructured = False
            logger.warning("unstructured library not available, falling back to basic loaders")
    
    def load_document(self, path: Path) -> tuple[str, list[dict]]:
        """Load document and extract text content and images"""
        ext = path.suffix.lower()
        
        if ext not in SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file type: {ext} for {path.name}")
            return "", []
        
        try:
            # For PDFs, always use pypdf first (it's reliable and already installed)
            if ext == '.pdf':
                return self._load_pdf(path)
            
            # Use unstructured library if available (best for complex formats except PDF)
            if self.has_unstructured and ext in {'.doc', '.docx', '.pptx', '.ppt', '.html', '.htm', '.eml', '.msg', '.epub', '.rtf', '.odt'}:
                text = self._load_with_unstructured(path)
                if text:  # If unstructured succeeds, use it
                    return text, []
                # Otherwise fall through to specialized loaders
            
            # Fallback to specialized loaders
            if ext in {'.docx', '.doc'}:
                return self._load_docx(path), []
            elif ext in {'.pptx', '.ppt'}:
                return self._load_pptx(path), []
            elif ext in {'.xlsx', '.xls'}:
                return self._load_xlsx(path), []
            elif ext in {'.html', '.htm', '.xml'}:
                return self._load_html(path), []
            elif ext in {'.csv', '.tsv'}:
                return self._load_csv(path), []
            elif ext in {'.json', '.ndjson'}:
                return self._load_json(path), []
            elif ext in {'.yaml', '.yml'}:
                return self._load_yaml(path), []
            else:
                # Plain text formats
                return self._load_text(path), []
                
        except Exception as e:
            logger.error(f"Error loading {path.name}: {e}")
            return "", []
    
    def _load_with_unstructured(self, path: Path) -> str:
        """Load document using unstructured library"""
        try:
            from unstructured.partition.auto import partition
            
            elements = partition(str(path))
            text_parts = []
            
            for element in elements:
                text = str(element).strip()
                if text:
                    text_parts.append(text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Unstructured loader failed for {path.name}: {e}")
            return ""
    
    def _load_pdf(self, path: Path) -> tuple[str, list[dict]]:
        """Load PDF document with images"""
        try:
            from pypdf import PdfReader
            import base64
            
            reader = PdfReader(str(path))
            text_parts = []
            images = []
            
            # Check if encrypted
            if reader.is_encrypted:
                try:
                    reader.decrypt('')
                except:
                    logger.warning(f"PDF {path.name} is encrypted")
                    return "", []
            
            for page_num, page in enumerate(reader.pages):
                # Extract text
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(text)
                
                # Extract images from /XObject resources
                try:
                    if '/Resources' in page and '/XObject' in page['/Resources']:
                        xObject = page['/Resources']['/XObject']
                        if hasattr(xObject, 'get_object'):
                            xObject = xObject.get_object()
                        
                        for obj_name in xObject:
                            obj = xObject[obj_name]
                            
                            if hasattr(obj, 'get_object'):
                                obj = obj.get_object()
                            
                            if '/Subtype' in obj and obj['/Subtype'] == '/Image':
                                try:
                                    # Get image dimensions
                                    width = int(obj.get('/Width', 0))
                                    height = int(obj.get('/Height', 0))
                                    
                                    # Skip very small images (likely decorative)
                                    if width < 50 or height < 50:
                                        continue
                                    
                                    # Get image data
                                    data = obj.get_data()
                                    
                                    # Only store reasonably sized images (< 1MB)
                                    if len(data) > 1000000:
                                        logger.debug(f"Skipping large image ({len(data)} bytes) from page {page_num}")
                                        continue
                                    
                                    # Encode to base64
                                    img_base64 = base64.b64encode(data).decode('utf-8')
                                    
                                    # Determine format
                                    img_format = 'jpeg'
                                    if '/Filter' in obj:
                                        filter_val = obj['/Filter']
                                        filter_name = str(filter_val).replace('/', '')
                                        if 'DCT' in filter_name:
                                            img_format = 'jpeg'
                                        elif 'Flate' in filter_name or 'PNG' in filter_name:
                                            img_format = 'png'
                                    
                                    images.append({
                                        'page': page_num + 1,
                                        'format': img_format,
                                        'width': width,
                                        'height': height,
                                        'data': img_base64,
                                        'text_position': len(text_parts) - 1 if text_parts else 0
                                    })
                                    
                                    logger.debug(f"Extracted image from page {page_num + 1}: {width}x{height} {img_format}")
                                    
                                except Exception as img_err:
                                    logger.debug(f"Could not extract image: {img_err}")
                                    continue
                except Exception as resource_err:
                    logger.debug(f"Could not access page resources: {resource_err}")
                    continue
            
            text_content = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(text_parts)} text blocks and {len(images)} images from {path.name}")
            
            return text_content, images
            
        except Exception as e:
            logger.error(f"PDF loader failed for {path.name}: {e}")
            return "", []
    
    def _load_docx(self, path: Path) -> str:
        """Load Word documents"""
        try:
            from docx import Document
            doc = Document(path)
            parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        parts.append(row_text)
            
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"DOCX loader failed for {path.name}: {e}")
            return self._load_text(path)  # Fallback to text
    
    def _load_pptx(self, path: Path) -> str:
        """Load PowerPoint presentations"""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
            
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"PPTX loader failed for {path.name}: {e}")
            return ""
    
    def _load_xlsx(self, path: Path) -> str:
        """Load Excel spreadsheets with markdown table formatting"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True)
            parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                parts.append(f"## Sheet: {sheet_name}\n")
                
                # Collect all rows
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                
                # Filter out completely empty rows
                non_empty_rows = [row for row in rows if any(cell is not None and str(cell).strip() for cell in row)]
                
                if not non_empty_rows:
                    continue
                
                # Determine max columns
                max_cols = max(len(row) for row in non_empty_rows)
                
                # Format as markdown table
                table_lines = []
                for row_idx, row in enumerate(non_empty_rows):
                    # Pad row to max columns
                    row_data = [str(cell).strip() if cell is not None else "" for cell in row]
                    row_data.extend([""] * (max_cols - len(row_data)))
                    
                    # Create table row
                    table_line = "| " + " | ".join(row_data) + " |"
                    table_lines.append(table_line)
                    
                    # Add separator after first row (header)
                    if row_idx == 0:
                        separator = "| " + " | ".join(["---"] * max_cols) + " |"
                        table_lines.append(separator)
                
                parts.append("\n".join(table_lines))
            
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"XLSX loader failed for {path.name}: {e}")
            return ""
    
    def _load_html(self, path: Path) -> str:
        """Load HTML/XML documents with markdown table formatting"""
        try:
            from bs4 import BeautifulSoup
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Convert tables to markdown before extracting text
            for table in soup.find_all('table'):
                markdown_table = self._table_to_markdown(table)
                # Replace table with markdown version
                table.replace_with(BeautifulSoup(f"\n\n{markdown_table}\n\n", 'html.parser'))
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        except Exception as e:
            logger.error(f"HTML loader failed for {path.name}: {e}")
            return ""
    
    def _table_to_markdown(self, table) -> str:
        """Convert HTML table to markdown format"""
        rows = []
        
        # Process header rows
        headers = []
        for thead in table.find_all('thead'):
            for tr in thead.find_all('tr'):
                header_row = []
                for th in tr.find_all(['th', 'td']):
                    header_row.append(th.get_text(strip=True))
                if header_row:
                    headers.append(header_row)
        
        # If no thead, check first row in tbody
        if not headers:
            tbody = table.find('tbody')
            if tbody:
                first_tr = tbody.find('tr')
                if first_tr and first_tr.find('th'):
                    header_row = [th.get_text(strip=True) for th in first_tr.find_all(['th', 'td'])]
                    headers.append(header_row)
        
        # Process body rows
        body_rows = []
        for tbody in table.find_all('tbody'):
            for tr in tbody.find_all('tr'):
                # Skip if already processed as header
                if headers and tr == tbody.find('tr') and tr.find('th'):
                    continue
                row_data = []
                for td in tr.find_all(['td', 'th']):
                    row_data.append(td.get_text(strip=True))
                if row_data:
                    body_rows.append(row_data)
        
        # If no tbody, process all trs
        if not body_rows:
            for tr in table.find_all('tr'):
                # Skip header rows
                if headers and any(tr in thead.find_all('tr') for thead in table.find_all('thead')):
                    continue
                row_data = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
                if row_data:
                    body_rows.append(row_data)
        
        if not headers and not body_rows:
            return ""
        
        # Use first data row as header if no headers found
        if not headers and body_rows:
            headers = [body_rows[0]]
            body_rows = body_rows[1:]
        
        # Determine max columns
        all_rows = (headers if headers else []) + body_rows
        max_cols = max(len(row) for row in all_rows) if all_rows else 0
        
        if max_cols == 0:
            return ""
        
        # Build markdown table
        lines = []
        
        # Header
        if headers:
            for header_row in headers:
                # Pad to max columns
                header_row = list(header_row) + [""] * (max_cols - len(header_row))
                lines.append("| " + " | ".join(header_row) + " |")
            
            # Separator
            lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        
        # Body rows
        for row in body_rows:
            row = list(row) + [""] * (max_cols - len(row))
            lines.append("| " + " | ".join(row) + " |")
        
        return "\n".join(lines)
    
    def _load_csv(self, path: Path) -> str:
        """Load CSV/TSV files with markdown table formatting"""
        try:
            import csv
            
            delimiter = '\t' if path.suffix.lower() == '.tsv' else ','
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = list(reader)
            
            if not rows:
                return ""
            
            # Filter out empty rows
            non_empty_rows = [row for row in rows if any(cell.strip() for cell in row)]
            
            if not non_empty_rows:
                return ""
            
            # Determine max columns
            max_cols = max(len(row) for row in non_empty_rows)
            
            # Format as markdown table
            lines = []
            for row_idx, row in enumerate(non_empty_rows):
                # Pad row to max columns
                row_data = [cell.strip() for cell in row]
                row_data.extend([""] * (max_cols - len(row_data)))
                
                # Create table row
                lines.append("| " + " | ".join(row_data) + " |")
                
                # Add separator after first row (header)
                if row_idx == 0:
                    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
            
            return "\n".join(lines)
        except Exception as e:
            logger.error(f"CSV loader failed for {path.name}: {e}")
            return self._load_text(path)
    
    def _load_json(self, path: Path) -> str:
        """Load JSON/NDJSON files"""
        try:
            import json
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                if path.suffix.lower() == '.ndjson':
                    # Newline-delimited JSON
                    data = [json.loads(line) for line in f if line.strip()]
                else:
                    data = json.load(f)
            
            # Convert to readable text
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            logger.error(f"JSON loader failed for {path.name}: {e}")
            return self._load_text(path)
    
    def _load_yaml(self, path: Path) -> str:
        """Load YAML files"""
        try:
            import yaml
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                data = yaml.safe_load(f)
            
            # Convert to readable text
            import json
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            logger.error(f"YAML loader failed for {path.name}: {e}")
            return self._load_text(path)
    
    def _load_text(self, path: Path) -> str:
        """Load plain text files with encoding detection"""
        try:
            # Try UTF-8 first
            try:
                return path.read_text(encoding='utf-8')
            except UnicodeDecodeError:
                # Try other encodings
                for encoding in ['utf-16', 'utf-32', 'latin-1', 'cp1252']:
                    try:
                        return path.read_text(encoding=encoding)
                    except:
                        continue
                # Last resort: ignore errors
                return path.read_text(encoding='utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"Text loader failed for {path.name}: {e}")
            return ""
    
    def iter_documents(self, docs_dir: str):
        """Iterate over all supported documents in directory"""
        base = Path(docs_dir)
        
        for p in base.rglob("*"):
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
                # Skip files we don't want
                if 'dont_want' in p.name.lower() or p.name.startswith('.'):
                    continue
                yield p
    
    def get_supported_extensions(self) -> set:
        """Return set of supported file extensions"""
        return SUPPORTED_EXTENSIONS


# Singleton instance
_loader = None

def get_loader() -> UnstructuredDocumentLoader:
    """Get the singleton loader instance"""
    global _loader
    if _loader is None:
        _loader = UnstructuredDocumentLoader()
    return _loader


def load_document(path: Path) -> str:
    """Convenience function to load a document"""
    return get_loader().load_document(path)


def iter_docs(docs_dir: str):
    """Convenience function to iterate documents"""
    return get_loader().iter_documents(docs_dir)
